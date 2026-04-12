import os
import fitz  # PyMuPDF
from pptx import Presentation
from pathlib import Path

def extract_text_from_file(file_path: str) -> str:
    """Extracts text content from PDF, PPTX, or code files for evaluation."""
    if not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            return text
            
        elif ext in [".ppt", ".pptx"]:
            text = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
            
        elif ext in [".py", ".cpp", ".txt", ".js", ".java", ".c", ".h"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
                
        else:
            return f"[Unsupported file type: {ext}]"
            
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return f"[Error extracting text: {str(e)}]"
